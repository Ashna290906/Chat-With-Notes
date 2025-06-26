from typing import List, Union, BinaryIO, Optional
import fitz
import pandas as pd
from docx import Document
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
import logging

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

def extract_text(file: BinaryIO, filetype: str) -> str:
    """
    Extract text from various file formats.
    
    Args:
        file: File-like object containing the document
        filetype: Type of the file (pdf, docx, pptx, xlsx, xls)
        
    Returns:
        str: Extracted text from the document
        
    Raises:
        ValueError: If the file type is not supported
        Exception: For any other errors during text extraction
    """
    try:
        file.seek(0)  # Ensure we're at the start of the file
        
        if filetype == "pdf":
            try:
                # Open the PDF with error handling
                doc = fitz.open(stream=file.read(), filetype="pdf")
                # Extract text from each page
                text_parts = []
                for page in doc:
                    text = page.get_text()
                    if text.strip():  # Only add non-empty pages
                        text_parts.append(text)
                return "\n".join(text_parts)
            except Exception as e:
                logger.error(f"Error reading PDF: {str(e)}")
                raise
                
        elif filetype == "docx":
            try:
                doc = Document(file)
                return "\n\n".join(
                    p.text for p in doc.paragraphs 
                    if p.text.strip()
                )
            except Exception as e:
                logger.error(f"Error reading DOCX: {str(e)}")
                raise
                
        elif filetype == "pptx":
            try:
                prs = Presentation(file)
                return "\n".join(
                    shape.text for slide in prs.slides 
                    for shape in slide.shapes 
                    if hasattr(shape, "text") and shape.text.strip()
                )
            except Exception as e:
                logger.error(f"Error reading PPTX: {str(e)}")
                raise
                
        elif filetype in ["xlsx", "xls"]:
            try:
                excel_data = pd.ExcelFile(file)
                text_parts = []
                
                for sheet_name in excel_data.sheet_names:
                    df = pd.read_excel(excel_data, sheet_name=sheet_name)
                    # Only include non-empty dataframes
                    if not df.empty:
                        text_parts.append(f"--- Sheet: {sheet_name} ---")
                        text_parts.append(df.to_string())
                
                return "\n\n".join(text_parts) if text_parts else ""
                
            except Exception as e:
                logger.error(f"Error reading Excel file: {str(e)}")
                raise
                
        else:
            raise ValueError(f"Unsupported file type: {filetype}")
            
    except Exception as e:
        logger.error(f"Error in extract_text: {str(e)}")
        raise

def split_text(text: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> List[str]:
    """
    Split text into chunks using a recursive character text splitter.
    
    Args:
        text: The text to split
        chunk_size: Maximum size of each chunk (in characters)
        chunk_overlap: Number of characters to overlap between chunks
        
    Returns:
        List[str]: List of text chunks
    """
    try:
        if not text or not text.strip():
            return []
            
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len,
            separators=["\n\n", "\n", ". ", " ", ""]
        )
        
        return splitter.split_text(text)
        
    except Exception as e:
        logger.error(f"Error in split_text: {str(e)}")
        # Fallback to simple splitting if the smart splitter fails
        return [text[i:i + chunk_size] for i in range(0, len(text), chunk_size)]
